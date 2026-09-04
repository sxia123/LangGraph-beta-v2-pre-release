import io
import os

import docx
import pptx
import pytest

from server import create_direct_chat_graph
from src.agents.chart_pipeline import create_chart_pipeline_graph
from src.core.local_llm import LocalLLMClient, LocalLLMConfig


@pytest.fixture
def mock_client():
    os.environ["LLM_PROVIDER"] = "mock"
    return LocalLLMClient(LocalLLMConfig(provider="mock"))


@pytest.fixture
def mock_chart_graph(mock_client):
    return create_chart_pipeline_graph(mock_client)


@pytest.fixture
def mock_direct_graph(mock_client):
    return create_direct_chat_graph(mock_client)


def test_routing_with_foreign_file_on_disk(mock_chart_graph):
    """If input references an existing foreign spreadsheet file, route to the spreadsheet pipeline."""
    nodes = []
    final_output = ""
    for chunk in mock_chart_graph.stream(
        {
            "user_input": "Please decipher mock_enterprise_financial_q1_2026.xlsx and show charts",
            "current_step": "intake",
        }
    ):
        for node, upd in chunk.items():
            nodes.append(node)
            if "final_response" in upd:
                final_output = upd["final_response"]

    assert "foreign_file_router" in nodes
    assert "spreadsheet_specialist" in nodes
    assert "spreadsheet_verify" in nodes
    # Must NOT run the original specialist web search pipeline
    assert "tier0_checks" not in nodes
    assert "tier05_web_verify" not in nodes
    # Check that output contains deciphered tables and Mermaid chart
    assert "Executive Summary" in final_output
    assert "```mermaid" in final_output


def test_routing_with_foreign_file_attachment_in_state(mock_chart_graph):
    """If input includes a foreign file in state['files'], route to spreadsheet pipeline."""
    nodes = []
    for chunk in mock_chart_graph.stream(
        {
            "user_input": "Analyze company sales data",
            "files": [{"filename": "company_sales_q1_2026.xlsx", "content": "company_sales_q1_2026.xlsx"}],
            "current_step": "intake",
        }
    ):
        for node in chunk.keys():
            nodes.append(node)

    assert "foreign_file_router" in nodes
    assert "spreadsheet_specialist" in nodes
    assert "spreadsheet_verify" in nodes
    assert "tier0_checks" not in nodes


def test_routing_with_document_file(mock_chart_graph):
    """If input includes a Word document (.docx), route to media specialist and generate diagram."""
    doc = docx.Document()
    doc.add_heading("Operational Framework 2026", 0)
    doc.add_paragraph("Policy guidelines for autonomous multi-agent verification systems.")
    buf = io.BytesIO()
    doc.save(buf)

    nodes = []
    final_output = ""
    for chunk in mock_chart_graph.stream(
        {
            "user_input": "Analyze the attached policy document and diagram the architecture",
            "files": [{"filename": "framework.docx", "content": buf.getvalue()}],
            "current_step": "intake",
        }
    ):
        for node, upd in chunk.items():
            nodes.append(node)
            if "final_response" in upd:
                final_output = upd["final_response"]

    assert "foreign_file_router" in nodes
    assert "spreadsheet_specialist" in nodes
    assert "spreadsheet_verify" in nodes
    assert "tier0_checks" not in nodes
    assert "```mermaid" in final_output


def test_routing_with_slideshow_file(mock_chart_graph):
    """If input includes a PowerPoint slideshow (.pptx), route to media specialist and generate roadmap."""
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Q1 Strategic Execution Roadmap"
    p_buf = io.BytesIO()
    prs.save(p_buf)

    nodes = []
    final_output = ""
    for chunk in mock_chart_graph.stream(
        {
            "user_input": "Review the attached presentation and construct a visual milestone timeline",
            "files": [{"filename": "roadmap.pptx", "content": p_buf.getvalue()}],
            "current_step": "intake",
        }
    ):
        for node, upd in chunk.items():
            nodes.append(node)
            if "final_response" in upd:
                final_output = upd["final_response"]

    assert "foreign_file_router" in nodes
    assert "spreadsheet_specialist" in nodes
    assert "spreadsheet_verify" in nodes
    assert "tier0_checks" not in nodes
    assert "```mermaid" in final_output


def test_routing_with_photo_image(mock_chart_graph):
    """If input includes photos/images, route to media specialist with vision capability."""
    nodes = []
    final_output = ""
    sample_img = "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
    for chunk in mock_chart_graph.stream(
        {
            "user_input": "Analyze this architecture screenshot and create a workflow chart",
            "images": [sample_img],
            "current_step": "intake",
        }
    ):
        for node, upd in chunk.items():
            nodes.append(node)
            if "final_response" in upd:
                final_output = upd["final_response"]

    assert "foreign_file_router" in nodes
    assert "spreadsheet_specialist" in nodes
    assert "spreadsheet_verify" in nodes
    assert "tier0_checks" not in nodes
    assert "```mermaid" in final_output


def test_direct_vision_intake(mock_direct_graph):
    """Direct vision pipeline intakes photos and documents cleanly."""
    sample_img = "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
    result = mock_direct_graph.invoke(
        {
            "messages": [
                {
                    "role": "user",
                    "content": "Please analyze this attached photo",
                    "images": [sample_img],
                }
            ],
            "images": [sample_img],
            "user_input": "Please analyze this attached photo",
        },
        config={"configurable": {"thread_id": "test_direct_vision"}},
    )
    assert result is not None
    response = result.get("final_response") or (result.get("messages")[-1].get("content") if result.get("messages") else "")
    assert "Visual Analysis" in response


def test_routing_without_foreign_file_keeps_original_pipeline(mock_chart_graph):
    """If input does NOT include any foreign file or photo, keep with the original pipeline."""
    nodes = []
    for chunk in mock_chart_graph.stream(
        {
            "user_input": "Explain how transformer attention heads work in modern LLMs",
            "current_step": "intake",
        }
    ):
        for node in chunk.keys():
            nodes.append(node)

    assert "foreign_file_router" in nodes
    # Must follow original pipeline
    assert "specialist" in nodes
    assert "tier0_checks" in nodes
    assert "tier05_web_verify" in nodes
    assert "tier1_verify" in nodes
    assert "revisions" in nodes
    assert "tier2_verify" in nodes
    # Must NOT route to spreadsheet pipeline
    assert "spreadsheet_specialist" not in nodes
    assert "spreadsheet_verify" not in nodes
