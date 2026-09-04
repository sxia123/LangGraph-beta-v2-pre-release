"""Fabricates realistic PPTX and XLSX sample files for testing AI intake and charting capabilities."""

import os

import docx
import openpyxl
import pptx
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


def fabricate_xlsx(filepath: str) -> None:
    wb = openpyxl.Workbook()

    # Define common styles
    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
    accent_fill = PatternFill(start_color="D9E1F2", end_color="D9E1F2", fill_type="solid")
    thin_border = Border(
        left=Side(style="thin", color="D9D9D9"),
        right=Side(style="thin", color="D9D9D9"),
        top=Side(style="thin", color="D9D9D9"),
        bottom=Side(style="thin", color="D9D9D9"),
    )

    # -------------------------------------------------------------
    # Sheet 1: Revenue Summary
    # -------------------------------------------------------------
    ws1 = wb.active
    ws1.title = "Revenue Summary"
    ws1.views.sheetView[0].showGridLines = True

    headers1 = ["Region", "Target Revenue ($)", "Actual Revenue ($)", "Variance ($)", "Attainment (%)", "Status"]
    data1 = [
        ["North America Enterprise", 4500000, 5250000, 750000, 1.167, "Exceeded"],
        ["North America Mid-Market", 2200000, 2480000, 280000, 1.127, "Exceeded"],
        ["Europe / UK Financial AI", 3800000, 3650000, -150000, 0.961, "On Track"],
        ["Asia Pacific Agents SDK", 1900000, 2420000, 520000, 1.274, "Exceeded"],
        ["Latin America Emerging", 800000, 890000, 90000, 1.113, "Exceeded"],
        ["Global Strategic Accounts", 5000000, 5950000, 950000, 1.190, "Exceeded"],
    ]

    ws1.append(headers1)
    for row in data1:
        ws1.append(row)

    # Totals Row
    ws1.append(["Total Consolidated", 18200000, 20640000, 2440000, 1.134, "Portfolio Overachieved"])

    # Style Header Row
    for col_idx in range(1, len(headers1) + 1):
        cell = ws1.cell(row=1, column=col_idx)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center", vertical="center")

    # Style Data Rows
    for r in range(2, len(data1) + 2):
        for c in range(1, len(headers1) + 1):
            cell = ws1.cell(row=r, column=c)
            cell.border = thin_border
            if c in [2, 3, 4]:
                cell.number_format = "$#,##0"
                cell.alignment = Alignment(horizontal="right")
            elif c == 5:
                cell.number_format = "0.0%"
                cell.alignment = Alignment(horizontal="right")
            elif c == 6:
                cell.alignment = Alignment(horizontal="center")

    # Style Total Row
    tot_row = len(data1) + 2
    for c in range(1, len(headers1) + 1):
        cell = ws1.cell(row=tot_row, column=c)
        cell.font = Font(name="Calibri", size=11, bold=True)
        cell.fill = accent_fill
        cell.border = thin_border
        if c in [2, 3, 4]:
            cell.number_format = "$#,##0"
            cell.alignment = Alignment(horizontal="right")
        elif c == 5:
            cell.number_format = "0.0%"
            cell.alignment = Alignment(horizontal="right")

    # -------------------------------------------------------------
    # Sheet 2: Product Performance
    # -------------------------------------------------------------
    ws2 = wb.create_sheet(title="Product Performance")
    ws2.views.sheetView[0].showGridLines = True

    headers2 = ["Product Line", "Active Enterprise Orgs", "ARR ($)", "Gross Margin (%)", "Net Expansion Rate (%)", "Tier"]
    data2 = [
        ["Autonomous Agents SDK", 620, 8400000, 0.845, 1.34, "Tier 1 Flagship"],
        ["Qwen Local Inference Engine", 410, 5200000, 0.880, 1.28, "Tier 1 Flagship"],
        ["Enterprise Graph Verification", 295, 4150000, 0.795, 1.22, "Tier 2 Growth"],
        ["Real-time Tool Orchestrator", 180, 2890000, 0.810, 1.18, "Tier 2 Growth"],
    ]

    ws2.append(headers2)
    for row in data2:
        ws2.append(row)

    for col_idx in range(1, len(headers2) + 1):
        cell = ws2.cell(row=1, column=col_idx)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center", vertical="center")

    for r in range(2, len(data2) + 2):
        for c in range(1, len(headers2) + 1):
            cell = ws2.cell(row=r, column=c)
            cell.border = thin_border
            if c in [2]:
                cell.number_format = "#,##0"
                cell.alignment = Alignment(horizontal="right")
            elif c in [3]:
                cell.number_format = "$#,##0"
                cell.alignment = Alignment(horizontal="right")
            elif c in [4, 5]:
                cell.number_format = "0.0%"
                cell.alignment = Alignment(horizontal="right")
            elif c == 6:
                cell.alignment = Alignment(horizontal="center")

    # Adjust column widths automatically
    for ws in [ws1, ws2]:
        for col in ws.columns:
            max_len = max(len(str(cell.value or "")) for cell in col)
            col_letter = get_column_letter(col[0].column)
            ws.column_dimensions[col_letter].width = max(max_len + 4, 14)

    wb.save(filepath)
    print(f"Created Excel spreadsheet: {filepath}")


def fabricate_pptx(filepath: str) -> None:
    prs = pptx.Presentation()
    # 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # -------------------------------------------------------------
    # Slide 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)

    # Title box
    tx_box1 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.5))
    tf1 = tx_box1.text_frame
    tf1.word_wrap = True

    p1 = tf1.paragraphs[0]
    p1.text = "Enterprise Multi-Agent AI Platform"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    p2 = tf1.add_paragraph()
    p2.text = "Q1 2026 Strategic Architecture, Verification Gates & Roadmap"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0x59, 0x59, 0x59)

    p3 = tf1.add_paragraph()
    p3.text = "Prepared for Executive Leadership & Engineering Steering Committee"
    p3.font.size = Pt(14)
    p3.font.italic = True
    p3.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)

    # Speaker notes for slide 1
    notes1 = slide1.notes_slide.notes_text_frame
    notes1.text = "Welcome executives. Today we present the multi-tier verification pipeline, foreign file intake metrics, and our 2026 expansion roadmap."

    # -------------------------------------------------------------
    # Slide 2: Strategic Pillars & Architecture
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)

    tx_box2 = slide2.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.0))
    tf2 = tx_box2.text_frame
    p_head2 = tf2.paragraphs[0]
    p_head2.text = "Strategic Architecture & Governance Pillars"
    p_head2.font.size = Pt(32)
    p_head2.font.bold = True
    p_head2.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    # Content box
    body_box2 = slide2.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.8))
    b_tf2 = body_box2.text_frame
    b_tf2.word_wrap = True

    points = [
        ("Multi-Modal Media Intake", "Unified parser handles Spreadsheets (.xlsx, .csv), Documents (.pdf, .docx), Slideshows (.pptx), and Photos (.png, .jpg)."),
        ("Deterministic 2.5-Stage Verification", "Enforces Tier 0 integrity checks, Tier 0.5 web verification, Tier 1 multi-agent audit, and final sanity gates."),
        ("Multi-Model Orchestration", "Local Qwen 27B / Qwen 2.5-VL for fast sub-second reasoning coupled with frontier model escalation pathways."),
        ("Persistent SQLite Audit Trails", "Every run, thought stream, file checkpoint, and human approval is cryptographically indexed in SQLite."),
    ]

    for idx, (title, desc) in enumerate(points):
        p_item = b_tf2.add_paragraph() if idx > 0 else b_tf2.paragraphs[0]
        p_item.text = f"- {title}: {desc}"
        p_item.font.size = Pt(18)
        p_item.space_after = Pt(14)

    notes2 = slide2.notes_slide.notes_text_frame
    notes2.text = "Highlight the deterministic verification gates and how foreign files are segregated safely into specialized deciphering workflows."

    # -------------------------------------------------------------
    # Slide 3: Milestone Timeline & Operational Metrics Table
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)

    tx_box3 = slide3.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.0))
    tf3 = tx_box3.text_frame
    p_head3 = tf3.paragraphs[0]
    p_head3.text = "Operational Milestones & Resource Commitments"
    p_head3.font.size = Pt(32)
    p_head3.font.bold = True
    p_head3.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    # Table on Slide 3
    rows = 5
    cols = 4
    left = Inches(1.0)
    top = Inches(2.2)
    width = Inches(11.3)
    height = Inches(3.8)

    table_shape = slide3.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table

    table_data = [
        ["Phase", "Milestone Target", "Target Delivery", "Status"],
        ["Phase 1", "Foreign File Intake (Excel, Word, Slides, Photos)", "Feb 2026", "Completed (Production)"],
        ["Phase 2", "Interactive Mermaid Visual Charting Engine", "Mar 2026", "Completed (Production)"],
        ["Phase 3", "Distributed GPU Cluster Scaling & Qwen-VL", "Apr 2026", "On Schedule"],
        ["Phase 4", "Autonomous Multi-Agent Enterprise Consensus", "May 2026", "Planned"],
    ]

    for r_idx, row in enumerate(table_data):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(15)
            if r_idx == 0:
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
            else:
                p.font.color.rgb = RGBColor(0x26, 0x26, 0x26)
                if r_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF4, 0xF8)

    notes3 = slide3.notes_slide.notes_text_frame
    notes3.text = "Emphasize that Phase 1 and Phase 2 are already completed and passing 100% test coverage."

    prs.save(filepath)
    print(f"Created PowerPoint presentation: {filepath}")


def fabricate_docx(filepath: str) -> None:
    """Fabricates a realistic enterprise Word document for testing Docling structured parsing."""
    doc = docx.Document()

    # Document Header
    doc.add_heading("Enterprise Multi-Agent AI Platform", 0)
    subtitle = doc.add_paragraph("Q1 2026 Technical Architecture, Media Ingestion & IBM Docling Verification Specification")
    subtitle.italic = True

    # Metadata Block
    meta = doc.add_paragraph()
    meta.add_run("Classification: ").bold = True
    meta.add_run("Internal Engineering / Strict Confidential\n")
    meta.add_run("Author: ").bold = True
    meta.add_run("Principal Systems Architect\n")
    meta.add_run("Release Version: ").bold = True
    meta.add_run("v2.4.0-Production\n")
    meta.add_run("Verification Engine: ").bold = True
    meta.add_run("IBM Docling 2.121.0\n")

    # Section 1
    doc.add_heading("1. Executive Summary & Architecture", level=1)
    doc.add_paragraph(
        "The Enterprise Multi-Agent AI Platform provides an end-to-end framework for autonomous agent orchestration, "
        "deterministic verification gates, and heterogeneous foreign document intake. In Q1 2026, the architecture "
        "integrates IBM Docling as the primary structured parser to decipher rich documents including Word documents (.docx), "
        "PDFs, and presentations into high-fidelity structured Markdown."
    )
    doc.add_paragraph(
        "By enforcing strict Tier 0 integrity checks, Tier 0.5 web verification, Tier 1 multi-agent audit, and final sanity gates, "
        "the platform ensures 100% verifiable outputs with full cryptographic SQLite audit trails."
    )

    # Section 2
    doc.add_heading("2. Media Intake Subsystems & Parser Specifications", level=1)
    doc.add_paragraph(
        "The following operational matrix details the intake subsystems, target media types, parser engines, and latency SLAs:"
    )

    subsystems = [
        ("Spreadsheet Ingestion", ".xlsx, .xlsm, .csv", "OpenPyXL Engine", "< 120ms", "Tier 0 Audit"),
        ("Presentation Parser", ".pptx, .ppt", "Python-PPTX Engine", "< 250ms", "Tier 0.5 Gate"),
        ("Structured Document Parser", ".docx, .doc, .pdf", "IBM Docling Engine", "< 180ms", "Tier 1 Gate"),
        ("Visual Media Ingestion", ".png, .jpg, .webp", "Qwen 2.5-VL Vision", "< 450ms", "Tier 2 Audit"),
        ("Sandboxed Tool Runtime", "Python, Terminal, Git", "Sandboxed Execution", "< 800ms", "Tier 2.5 Sanity"),
    ]

    table = doc.add_table(rows=1, cols=5)
    hdr_cells = table.rows[0].cells
    hdr_titles = ["Subsystem", "Formats", "Parser Engine", "Latency SLA", "Verification Gate"]
    for i, title in enumerate(hdr_titles):
        hdr_cells[i].text = title
        hdr_cells[i].paragraphs[0].runs[0].bold = True

    for item in subsystems:
        row_cells = table.add_row().cells
        for idx, val in enumerate(item):
            row_cells[idx].text = val

    # Section 3
    doc.add_heading("3. Strategic Roadmap & Milestones", level=1)
    doc.add_paragraph("Key strategic milestones for the current release cycle include:")
    milestones = [
        "Milestone 1: Complete integration of IBM Docling 2.x for high-precision table extraction and Markdown generation.",
        "Milestone 2: Browser-based Playwright automated UI verification across Edge and Chromium headless runtimes.",
        "Milestone 3: Local LLM inference acceleration via Qwen 3.5 9B with graceful fallbacks.",
        "Milestone 4: Interactive Mermaid diagram synthesis for document workflows and pipeline topologies.",
    ]
    for m in milestones:
        doc.add_paragraph(m, style="List Bullet")

    # Section 4
    doc.add_heading("4. Verification & Operational Sign-off", level=1)
    doc.add_paragraph(
        "All document ingestion pipelines and automated browser verification tests have completed with zero regressions. "
        "System status: CERTIFIED OPERATIONAL."
    )

    doc.save(filepath)
    print(f"Created Enterprise Word document: {filepath}")


if __name__ == "__main__":
    xlsx_path = os.path.abspath("test_global_sales_q1_2026.xlsx")
    pptx_path = os.path.abspath("test_ai_product_roadmap_2026.pptx")
    docx_path = os.path.abspath("test_enterprise_document_2026.docx")

    fabricate_xlsx(xlsx_path)
    fabricate_pptx(pptx_path)
    fabricate_docx(docx_path)

