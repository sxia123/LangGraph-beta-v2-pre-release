"""Document and media parser engine for photos, documents (PDF, DOCX), and slideshows (PPTX).

Supports both raw binary content, file paths, and base64 data URIs.
Extracts structured Markdown text, tables, slide progressions, and metadata.
"""

import base64
import io
import logging
import os
import re
from typing import Any, Dict, List, Optional, Union

import docx
import pptx
import pypdfium2

from src.core.spreadsheet_parser import decipher_spreadsheet

logger = logging.getLogger("document_parser")


def _decode_source_to_bytes(source: Union[str, bytes]) -> tuple[Optional[bytes], Optional[str]]:
    """Converts a string (data URI, file path, or base64) or bytes into (raw_bytes, filepath)."""
    if isinstance(source, bytes):
        return source, None

    if not isinstance(source, str):
        return None, None

    # Check for base64 data URI (e.g. data:application/pdf;base64,...)
    if source.startswith("data:") and ";base64," in source:
        try:
            _, encoded = source.split(";base64,", 1)
            return base64.b64decode(encoded), None
        except Exception as err:
            logger.warning("Failed to decode base64 data URI: %s", err)
            return None, None

    # Check if string is an existing file path
    if os.path.isfile(source):
        try:
            with open(source, "rb") as f:
                return f.read(), source
        except Exception as err:
            logger.warning("Failed to read file path %s: %s", source, err)
            return None, source

    # Try raw base64 decode if it appears to be base64
    clean = source.strip()
    if len(clean) > 32 and re.match(r"^[A-Za-z0-9+/=\r\n]+$", clean):
        try:
            return base64.b64decode(clean), None
        except Exception:
            pass

    # Fallback to UTF-8 encoded bytes of text
    return source.encode("utf-8"), None


def parse_pdf(
    source: Union[str, bytes],
    filename: Optional[str] = None,
    max_pages: int = 50,
) -> Dict[str, Any]:
    """Extracts structured text from a PDF document using pypdfium2.

    Returns a dict with 'ok', 'type', 'filename', 'page_count', 'extracted_text', and 'summary'.
    """
    fn = filename or "document.pdf"
    raw_bytes, path = _decode_source_to_bytes(source)
    if not raw_bytes:
        return {
            "ok": False,
            "type": "document",
            "filename": fn,
            "error": "Unable to read PDF source data.",
            "extracted_text": "",
        }

    try:
        pdf = pypdfium2.PdfDocument(path if path else io.BytesIO(raw_bytes))
        total_pages = len(pdf)
        pages_to_read = min(total_pages, max_pages)

        page_texts: List[str] = []
        for i in range(pages_to_read):
            page = pdf[i]
            text_page = page.get_textpage()
            text = text_page.get_text_range() or ""
            text_cleaned = text.strip()
            if text_cleaned:
                page_texts.append(f"### Page {i + 1}\n{text_cleaned}")

        combined_body = "\n\n".join(page_texts) if page_texts else "[No readable text found on document pages.]"
        omitted_note = f"\n\n*(Note: {total_pages - pages_to_read} additional pages omitted)*" if total_pages > pages_to_read else ""

        full_content = (
            f"## Deciphered PDF Document: {fn}\n"
            f"- **Total Pages**: {total_pages}\n"
            f"- **Pages Processed**: {pages_to_read}\n\n"
            f"{combined_body}{omitted_note}"
        )

        return {
            "ok": True,
            "type": "document",
            "format": "pdf",
            "filename": fn,
            "page_count": total_pages,
            "extracted_text": full_content,
            "deciphered_context": full_content,
            "summary": f"**PDF Document**: `{fn}` ({total_pages} page(s), {len(page_texts)} non-empty)",
        }
    except Exception as err:
        logger.error("Error parsing PDF %s: %s", fn, err)
        return {
            "ok": False,
            "type": "document",
            "filename": fn,
            "error": str(err),
            "extracted_text": f"[Error reading PDF '{fn}': {err}]",
        }


def parse_docx(
    source: Union[str, bytes],
    filename: Optional[str] = None,
) -> Dict[str, Any]:
    """Extracts paragraphs, headings, and embedded tables from a Word document (.docx).

    Returns a dict with 'ok', 'type', 'filename', 'extracted_text', and 'summary'.
    """
    fn = filename or "document.docx"
    raw_bytes, path = _decode_source_to_bytes(source)
    if not raw_bytes:
        return {
            "ok": False,
            "type": "document",
            "filename": fn,
            "error": "Unable to read DOCX source data.",
            "extracted_text": "",
        }

    try:
        stream = io.BytesIO(raw_bytes) if not path else path
        doc = docx.Document(stream)

        sections: List[str] = []

        # 1. Paragraphs
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        if paragraphs:
            sections.append("### Document Body\n" + "\n\n".join(paragraphs))

        # 2. Tables
        if doc.tables:
            table_md_blocks: List[str] = []
            for t_idx, table in enumerate(doc.tables):
                rows = []
                for row in table.rows:
                    cells = [c.text.replace("\n", " ").strip() for c in row.cells]
                    # Deduplicate repeated cell texts from merged cells
                    cleaned_cells = [str(c).replace("|", "\\|") for c in cells]
                    rows.append(cleaned_cells)

                if rows:
                    header = rows[0]
                    md = [
                        f"#### Table {t_idx + 1}",
                        "| " + " | ".join(header) + " |",
                        "| " + " | ".join(["---"] * max(1, len(header))) + " |",
                    ]
                    for r in rows[1:100]:
                        md.append("| " + " | ".join(r) + " |")
                    if len(rows) > 101:
                        md.append(f"\n*... ({len(rows) - 101} additional rows omitted)*")
                    table_md_blocks.append("\n".join(md))

            if table_md_blocks:
                sections.append("### Embedded Tables\n" + "\n\n".join(table_md_blocks))

        content_body = "\n\n".join(sections) if sections else "[Empty Word Document]"
        full_content = (
            f"## Deciphered Word Document: {fn}\n"
            f"- **Paragraphs**: {len(paragraphs)}\n"
            f"- **Tables**: {len(doc.tables)}\n\n"
            f"{content_body}"
        )

        return {
            "ok": True,
            "type": "document",
            "format": "docx",
            "filename": fn,
            "paragraphs_count": len(paragraphs),
            "tables_count": len(doc.tables),
            "extracted_text": full_content,
            "deciphered_context": full_content,
            "summary": f"**Word Document**: `{fn}` ({len(paragraphs)} paragraph(s), {len(doc.tables)} table(s))",
        }
    except Exception as err:
        logger.error("Error parsing DOCX %s: %s", fn, err)
        return {
            "ok": False,
            "type": "document",
            "filename": fn,
            "error": str(err),
            "extracted_text": f"[Error reading Word document '{fn}': {err}]",
        }


def parse_slideshow(
    source: Union[str, bytes],
    filename: Optional[str] = None,
    max_slides: int = 50,
) -> Dict[str, Any]:
    """Extracts titles, bullets, tables, and speaker notes from a presentation (.pptx).

    Returns a dict with 'ok', 'type', 'filename', 'slide_count', 'extracted_text', and 'summary'.
    """
    fn = filename or "presentation.pptx"
    raw_bytes, path = _decode_source_to_bytes(source)
    if not raw_bytes:
        return {
            "ok": False,
            "type": "slideshow",
            "filename": fn,
            "error": "Unable to read PPTX source data.",
            "extracted_text": "",
        }

    try:
        stream = io.BytesIO(raw_bytes) if not path else path
        prs = pptx.Presentation(stream)
        total_slides = len(prs.slides)
        slides_to_read = min(total_slides, max_slides)

        slide_blocks: List[str] = []
        for idx in range(slides_to_read):
            slide = prs.slides[idx]
            title_text = ""
            if slide.shapes.title and slide.shapes.title.text:
                title_text = slide.shapes.title.text.strip()

            bullets: List[str] = []
            tables_md: List[str] = []

            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        txt = p.text.strip()
                        if txt:
                            bullets.append(f"- {txt}")
                elif shape.has_table:
                    tbl = shape.table
                    tbl_rows = []
                    for row in tbl.rows:
                        cells = [c.text.replace("\n", " ").strip() for c in row.cells]
                        tbl_rows.append([str(c).replace("|", "\\|") for c in cells])
                    if tbl_rows:
                        h = tbl_rows[0]
                        tbl_str = [
                            "| " + " | ".join(h) + " |",
                            "| " + " | ".join(["---"] * max(1, len(h))) + " |",
                        ]
                        for r in tbl_rows[1:20]:
                            tbl_str.append("| " + " | ".join(r) + " |")
                        tables_md.append("\n".join(tbl_str))

            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                n_text = slide.notes_slide.notes_text_frame.text.strip()
                if n_text:
                    notes_text = f"\n*Speaker Notes*: {n_text}"

            slide_header = f"### Slide {idx + 1}: {title_text or '(Untitled Slide)'}"
            slide_parts = [slide_header]
            if bullets:
                slide_parts.append("\n".join(bullets))
            if tables_md:
                slide_parts.append("\n\n".join(tables_md))
            if notes_text:
                slide_parts.append(notes_text)

            slide_blocks.append("\n".join(slide_parts))

        content_body = "\n\n---\n\n".join(slide_blocks) if slide_blocks else "[Empty Presentation]"
        omitted_note = f"\n\n*(Note: {total_slides - slides_to_read} additional slides omitted)*" if total_slides > slides_to_read else ""

        full_content = (
            f"## Deciphered Slideshow: {fn}\n"
            f"- **Total Slides**: {total_slides}\n"
            f"- **Slides Processed**: {slides_to_read}\n\n"
            f"{content_body}{omitted_note}"
        )

        return {
            "ok": True,
            "type": "slideshow",
            "format": "pptx",
            "filename": fn,
            "slide_count": total_slides,
            "extracted_text": full_content,
            "deciphered_context": full_content,
            "summary": f"**PowerPoint Slideshow**: `{fn}` ({total_slides} slide(s))",
        }
    except Exception as err:
        logger.error("Error parsing PPTX %s: %s", fn, err)
        return {
            "ok": False,
            "type": "slideshow",
            "filename": fn,
            "error": str(err),
            "extracted_text": f"[Error reading PowerPoint slideshow '{fn}': {err}]",
        }


def decipher_media_file(
    source: Union[str, bytes],
    filename: Optional[str] = None,
) -> Dict[str, Any]:
    """Unified dispatcher that deciphers any supported foreign media file.

    Handles spreadsheets (.xlsx, .csv, etc.), documents (.pdf, .docx),
    slideshows (.pptx), and images (.png, .jpg, etc.).
    """
    fn = filename or "attached_file"
    fn_lower = fn.lower()

    # 1. Spreadsheets
    if fn_lower.endswith((".xlsx", ".xls", ".xlsm", ".csv", ".tsv")):
        res = decipher_spreadsheet(source, filename=fn)
        if res.get("ok"):
            res["type"] = "spreadsheet"
        return res

    # 2. PDF Documents
    if fn_lower.endswith(".pdf"):
        return parse_pdf(source, filename=fn)

    # 3. Word Documents
    if fn_lower.endswith((".docx", ".doc")):
        return parse_docx(source, filename=fn)

    # 4. Presentations / Slideshows
    if fn_lower.endswith((".pptx", ".ppt")):
        return parse_slideshow(source, filename=fn)

    # 5. Photos / Images
    if fn_lower.endswith((".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".svg")) or (
        isinstance(source, str) and source.startswith("data:image/")
    ):
        return {
            "ok": True,
            "type": "image",
            "filename": fn,
            "summary": f"**Visual Image**: `{fn}`",
            "deciphered_context": f"## Visual Image Asset: {fn}\nAttached image for visual inspection and architectural charting.",
        }

    # 6. Generic plain text or markdown fallback
    raw_bytes, _ = _decode_source_to_bytes(source)
    if raw_bytes:
        try:
            txt = raw_bytes.decode("utf-8")
        except Exception:
            txt = raw_bytes.decode("latin-1", errors="replace")
        return {
            "ok": True,
            "type": "document",
            "filename": fn,
            "extracted_text": txt,
            "deciphered_context": f"## Document Attachment: {fn}\n\n{txt}",
            "summary": f"**Document**: `{fn}`",
        }

    return {
        "ok": False,
        "type": "unknown",
        "filename": fn,
        "error": f"Unsupported or unreadable media format for '{fn}'.",
    }

